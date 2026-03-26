import os
from pptx import Presentation
from pptx.util import Inches, Pt

def markdown_to_pptx_v2(md_path, pptx_path, gif_path=None):
    prs = Presentation()
    
    # 讀取 MD
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()

    # 簡單的標題分割邏輯
    sections = content.split('## ')
    
    # 第一頁：大標題
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Isochrone Design Logic"
    subtitle.text = "Thesis Presentation - El Prat de Llobregat"

    for section in sections[1:]:
        lines = section.split('\n')
        title_text = lines[0].strip()
        body_text = "\n".join([line.strip() for line in lines[1:] if line.strip()])
        
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title_text
        
        tf = slide.placeholders[1].text_frame
        tf.text = body_text[:1000] # 避免過長
        
    # 最後一頁：動畫
    if gif_path and os.path.exists(gif_path):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # 加入文字框標題
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        txBox.text_frame.text = "Final Isochrone Animation"
        # 加入圖片
        slide.shapes.add_picture(gif_path, Inches(1), Inches(1.5), height=Inches(5))

    prs.save(pptx_path)
    print(f"PPTX saved to {pptx_path}")

if __name__ == "__main__":
    md_file = "/media/chi/ai_/Thesis_Oasis/Map_Isochrone_Data/Isochrone_Design_Logic_EN.md"
    pptx_file = "/media/chi/ai_/Thesis_Oasis/Map_Isochrone_Data/Thesis_Presentation_V2.pptx"
    gif_file = "/media/chi/ai_/Thesis_Oasis/Map_Isochrone_Data/thesis_smappen_fusion.gif"
    markdown_to_pptx_v2(md_file, pptx_file, gif_file)
