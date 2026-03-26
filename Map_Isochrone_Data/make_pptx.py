import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def markdown_to_pptx(md_path, pptx_path, gif_path=None):
    prs = Presentation()
    
    # Read MD content
    with open(md_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()

    current_slide = None
    
    for line in lines:
        line = line.strip()
        if not line:
            continue
            
        # H1 or H2 as Slide Title
        if line.startswith('# ') or line.startswith('## '):
            title_text = line.lstrip('#').strip()
            slide_layout = prs.slide_layouts[1] # Title and Content
            current_slide = prs.slides.add_slide(slide_layout)
            current_slide.shapes.title.text = title_text
            
        # Bullet points or text
        elif current_slide:
            body_shape = current_slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            p = tf.add_paragraph()
            p.text = line.lstrip('*').lstrip('-').strip()
            p.level = 0 if line.startswith(('*', '-')) else 0
            
    # Add a final slide for the Animation
    if gif_path and os.path.exists(gif_path):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        title = slide.shapes.title if slide.shapes.title else None
        
        # Add title manually to blank slide
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(1))
        tf = txBox.text_frame
        tf.text = "Final Animation: Isochrone Expansion"
        
        # Add the GIF (Note: PPTX might show static first frame depending on viewer)
        slide.shapes.add_picture(gif_path, Inches(1), Inches(1.5), height=Inches(5.5))

    prs.save(pptx_path)
    print(f"PPTX saved to {pptx_path}")

if __name__ == "__main__":
    md_file = "/media/chi/ai_/Thesis_Oasis/Map_Isochrone_Data/Isochrone_Design_Logic_EN.md"
    pptx_file = "/media/chi/ai_/Thesis_Oasis/Map_Isochrone_Data/Thesis_Presentation.pptx"
    gif_file = "/media/chi/ai_/Thesis_Oasis/Map_Isochrone_Data/thesis_smappen_fusion.gif"
    
    markdown_to_pptx(md_file, pptx_file, gif_file)
